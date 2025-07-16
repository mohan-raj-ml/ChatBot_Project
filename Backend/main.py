# main.py
from fastapi import FastAPI, Request, HTTPException, UploadFile, Form, File
from fastapi.responses import JSONResponse, Response
from fastapi.middleware.cors import CORSMiddleware
from starlette.middleware.sessions import SessionMiddleware
from pydantic import BaseModel
from typing import Optional
import sqlite3
import requests
import shutil
import os
import io
from pptx import Presentation
import asyncio
import logging
import re
import time
from fastapi import BackgroundTasks
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain.llms.base import LLM
from langchain.chains.summarize import load_summarize_chain
from langchain.text_splitter import CharacterTextSplitter
from langchain.docstore.document import Document


import tiktoken

response_time_logger = logging.getLogger("response_time")
response_time_logger.setLevel(logging.INFO)

# Create file handler for the logger
rt_file_handler = logging.FileHandler("logs/response_times.log")
rt_file_handler.setLevel(logging.INFO)

# Optional: formatter for clean timestamps
formatter = logging.Formatter('%(asctime)s - %(message)s')
rt_file_handler.setFormatter(formatter)

# Add the handler to the logger
response_time_logger.addHandler(rt_file_handler)

MODEL_TOKEN_LIMITS = {
    "gemma:2b": 2048,
    "gemma:7b": 8192,
    "mistral": 8192,
    "mistral:7b": 8192,
    "llama3": 8192,
    "llama3:8b": 8192,
    "llama2": 4096,
    "phi": 2048,
    "tinyllama": 2048,
    "default": 4096
}

def calculate_response_token_budget(full_context: str, model: str) -> int:
    model = model.lower()
    total_limit = next((v for k, v in MODEL_TOKEN_LIMITS.items() if k in model), MODEL_TOKEN_LIMITS["default"])
    used_tokens = count_tokens(full_context, model)
    
    # Leave a tiny safety buffer (e.g. 20 tokens)
    available_for_response = max(total_limit - used_tokens - 100, 128)
    
    return available_for_response



def count_tokens(text: str, model: str = "gpt-3.5-turbo") -> int:
    try:
        encoding = tiktoken.encoding_for_model(model)
    except Exception:
        encoding = tiktoken.get_encoding("cl100k_base")
    return len(encoding.encode(text))

def trim_context(summary, retrieved, recent, file_section, prompt, max_tokens=3800):
    parts = {
        "prompt": prompt,
        "retrieved": retrieved,
        "summary": summary,   
        "recent": recent,
        "file_section": file_section
        
    }
    used_tokens = 0
    trimmed = {}

    for key, value in parts.items():
        tokens = count_tokens(value)
        if used_tokens + tokens <= max_tokens:
            trimmed[key] = value
            used_tokens += tokens
        else:
            trimmed[key] = f"...[Omitted {key} due to token limit]..."

    return trimmed

def remove_think_tags(text: str) -> str:
    return re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL).strip()

# Assuming vector_store.py is in the same directory and contains the updated functions
from vector_store import get_faiss_index, save_faiss_index, embedding_model as global_embedding_model

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

llm_cache = {}

def get_llm(model: str):
    if model not in llm_cache:
        llm_cache[model] = OllamaLLM(model=model)
    return llm_cache[model]


faiss_cache = {}

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000","*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

app.add_middleware(
    SessionMiddleware,
    secret_key="supersecretkey", # IMPORTANT: Use a strong, truly secret key in production
    max_age=60 * 60 * 24 * 4, # 4 days
    same_site="lax",
    session_cookie="session"
)

DB_NAME = 'users.db'

def get_db():
    conn = sqlite3.connect(DB_NAME, check_same_thread=False)
    conn.row_factory = sqlite3.Row # Allows accessing columns by name
    conn.execute("PRAGMA foreign_keys = ON") # Enforce foreign key constraints
    return conn

def init_db():
    with sqlite3.connect(DB_NAME) as db:
        db.execute("PRAGMA foreign_keys = ON")
        # Check if schema.sql exists and execute it
        if os.path.exists('schema.sql'):
            with open('schema.sql', 'r') as f:
                db.executescript(f.read())
            logger.info("Database schema initialized successfully.")
        else:
            logger.warning("schema.sql not found. Database schema might not be initialized.")

class OllamaLLM(LLM):
    model: str = "mistral" # Default model, can be overridden

    def _call(self, prompt: str, stop=None):
        try:
            res = requests.post("http://localhost:11434/api/generate", json={
                "model": self.model,
                "prompt": prompt,
                "stream": False
            }, timeout=300) # Added timeout for robustness
            res.raise_for_status() # Raise an exception for HTTP errors (4xx or 5xx)
            return res.json().get("response", "").strip() # .strip() response
        except requests.exceptions.RequestException as e:
            logger.error(f"Error communicating with Ollama server ({self.model}): {e}")
            raise HTTPException(status_code=503, detail=f"Failed to connect to Ollama model '{self.model}'. Please ensure it's running.")
        except Exception as e:
            logger.error(f"Unexpected error in OllamaLLM call: {e}")
            raise HTTPException(status_code=500, detail=f"LLM generation failed: {e}")

    @property
    def _identifying_params(self):
        return {"model": self.model}

    @property
    def _llm_type(self):
        return "ollama_custom"

class User(BaseModel):
    username: str
    email: Optional[str] = None
    password: str

class LoginRequest(BaseModel):
    identifier: str  # username or email
    password: str

class PromptRequest(BaseModel): # This model isn't used directly in /api/respond, but good for clarity
    prompt: str
    model: str
    chat_id: int

class ChatRequest(BaseModel):
    title: Optional[str] = "New Chat"

class RenameChatRequest(BaseModel):
    chat_id: int
    new_title: str

@app.post("/signup")
def signup(user: User):
    db = get_db()
    email = user.email.strip()
    username = user.username.strip()
    password = user.password.strip()

    # ✅ Restrict allowed domain
    allowed_domain = "@xoriant.com"
    if not email.endswith(allowed_domain):
        raise HTTPException(status_code=400, detail=f"Email must be from the {allowed_domain} domain.")

    try:
        db.execute(
            "INSERT INTO users (username, email, password) VALUES (?, ?, ?)",
            (username, email, password)
        )
        db.commit()
        logger.info(f"User {username} signed up successfully.")
        return {"success": True, "message": "Signup successful."}
    except sqlite3.IntegrityError as e:
        if "username" in str(e):
            raise HTTPException(status_code=409, detail="Username already exists.")
        elif "email" in str(e):
            raise HTTPException(status_code=409, detail="Email already exists.")
        else:
            raise HTTPException(status_code=500, detail="Internal error during signup.")



@app.post("/login")
async def login(login_req: LoginRequest, request: Request):
    db = get_db()
    identifier = login_req.identifier.strip()
    password = login_req.password.strip()

    user = db.execute(
        "SELECT * FROM users WHERE (username = ? OR email = ?) AND password = ?",
        (identifier, identifier, password)
    ).fetchone()

    if user:
        request.session['user'] = user["username"]
        logger.info(f"User {user['username']} logged in successfully.")
        return {"success": True, "message": "Login successful."}
    
    raise HTTPException(status_code=401, detail="Invalid credentials.")


@app.post("/logout")
async def logout(request: Request):
    user = request.session.get("user", "unknown")
    request.session.clear()
    logger.info(f"User {user} logged out successfully.")
    return {"success": True, "message": "Logged out successfully."}

@app.get("/")
def index(request: Request):
    user = request.session.get("user")
    if not user:
        return JSONResponse(status_code=401, content={"authenticated": False, "message": "User not logged in"})
    return {"authenticated": True, "user": user}

@app.get("/api/get_models")
def get_models():
    try:
        res = requests.get("http://localhost:11434/api/tags", timeout=10) # Added timeout
        res.raise_for_status()
        models = [m['name'] for m in res.json().get("models", [])]
        logger.info(f"Successfully retrieved Ollama models: {models}")
        return {"response": [models]}
    except requests.exceptions.RequestException as e:
        logger.error(f"Error fetching Ollama models: {e}")
        return {"response": [[]], "error": "Could not connect to Ollama server."}
    except Exception as e:
        logger.exception("Unexpected error fetching Ollama models.")
        return {"response": [[]], "error": f"An unexpected error occurred: {e}"}

@app.get("/api/list_chats")
def list_chats(request: Request):
    username = request.session.get("user")
    if not username:
        raise HTTPException(status_code=401, detail="Unauthorized")
    db = get_db()
    user = db.execute("SELECT id FROM users WHERE username = ?", (username,)).fetchone()
    if not user:
        logger.warning(f"User {username} not found in DB but in session. Session cleared.")
        request.session.clear() # Invalidate session if user isn't found
        raise HTTPException(status_code=401, detail="Unauthorized: User profile missing.")
    chats = db.execute("SELECT id, title, created_at FROM chats WHERE user_id = ? ORDER BY created_at DESC", (user["id"],)).fetchall()
    return {"chats": [dict(row) for row in chats]}

@app.get("/api/chat_history")
def chat_history(chat_id: int, request: Request):
    username = request.session.get("user")
    if not username:
        raise HTTPException(status_code=401, detail="Unauthorized")
    db = get_db()
    # Optional: Verify chat_id belongs to the user for security
    user_id = db.execute("SELECT id FROM users WHERE username = ?", (username,)).fetchone()["id"]
    chat_owner = db.execute("SELECT user_id FROM chats WHERE id = ?", (chat_id,)).fetchone()
    if not chat_owner or chat_owner["user_id"] != user_id:
        logger.warning(f"Unauthorized access attempt to chat_id {chat_id} by user {username}.")
        raise HTTPException(status_code=403, detail="Forbidden: Chat does not belong to user.")

    messages = db.execute("SELECT role, content FROM messages WHERE chat_id = ? ORDER BY timestamp ASC", (chat_id,)).fetchall()
    return [dict(row) for row in messages]

@app.get("/api/shared_chat_history")
def shared_chat_history(chat_id: int):
    # This endpoint is public for sharing, no auth check needed here.
    db = get_db()
    messages = db.execute("SELECT role, content FROM messages WHERE chat_id = ? ORDER BY timestamp ASC", (chat_id,)).fetchall()
    if not messages:
        raise HTTPException(status_code=404, detail="Chat not found or empty.")
    return [dict(row) for row in messages]


@app.post("/api/create_chat")
def create_chat(req: ChatRequest, request: Request):
    username = request.session.get("user")
    if not username:
        raise HTTPException(status_code=401, detail="Unauthorized")
    db = get_db()
    user = db.execute("SELECT id FROM users WHERE username = ?", (username,)).fetchone()
    if not user:
        logger.warning(f"User {username} not found when creating chat. Session cleared.")
        request.session.clear()
        raise HTTPException(status_code=401, detail="Unauthorized: User profile missing.")

    try:
        db.execute("INSERT INTO chats (user_id, title) VALUES (?, ?)", (user["id"], req.title))
        db.commit()
        chat_id = db.execute("SELECT last_insert_rowid()").fetchone()[0]
        # Initialize chat_memory and chat_summaries entries
        db.execute("INSERT INTO chat_memory (chat_id, memory) VALUES (?, ?)", (chat_id, ""))
        db.execute("INSERT INTO chat_summaries (chat_id, summary) VALUES (?, ?)", (chat_id, ""))
        db.commit()
        logger.info(f"New chat '{req.title}' created with ID {chat_id} for user {username}.")
        return {"success": True, "chat_id": chat_id, "title": req.title}
    except Exception as e:
        logger.exception(f"Error creating chat for user {username}.")
        db.rollback()
        raise HTTPException(status_code=500, detail="Failed to create chat.")


@app.post("/api/rename_chat")
def rename_chat(req: RenameChatRequest, request: Request):
    username = request.session.get("user")
    if not username:
        raise HTTPException(status_code=401, detail="Unauthorized")
    db = get_db()
    try:
        # Optional: Verify chat_id belongs to the user for security
        user_id = db.execute("SELECT id FROM users WHERE username = ?", (username,)).fetchone()["id"]
        chat_owner = db.execute("SELECT user_id FROM chats WHERE id = ?", (req.chat_id,)).fetchone()
        if not chat_owner or chat_owner["user_id"] != user_id:
            logger.warning(f"Unauthorized rename attempt to chat_id {req.chat_id} by user {username}.")
            raise HTTPException(status_code=403, detail="Forbidden: Chat does not belong to user.")

        db.execute("UPDATE chats SET title = ? WHERE id = ? AND user_id = ?", (req.new_title, req.chat_id, user_id))
        db.commit()
        logger.info(f"Chat {req.chat_id} renamed to '{req.new_title}' by user {username}.")
        return {"success": True}
    except Exception as e:
        logger.exception(f"Error renaming chat {req.chat_id} for user {username}.")
        db.rollback()
        raise HTTPException(status_code=500, detail="Failed to rename chat.")


@app.post("/api/respond")

async def respond(
    request: Request,
    background_tasks: BackgroundTasks,
    prompt: str = Form(...),
    model: str = Form(...),
    chat_id: int = Form(...),
    file: UploadFile = File(None)
):
    start = time.time()
    username = request.session.get("user")
    if not username:
        raise HTTPException(status_code=401, detail="Unauthorized")

    db = get_db()
    current_user_id = None
    try:
        current_user_id = db.execute("SELECT id FROM users WHERE username = ?", (username,)).fetchone()["id"]
        chat_owner = db.execute("SELECT user_id FROM chats WHERE id = ?", (chat_id,)).fetchone()
        if not chat_owner or chat_owner["user_id"] != current_user_id:
            logger.warning(f"Unauthorized access attempt to chat_id {chat_id} by user {username} during respond.")
            raise HTTPException(status_code=403, detail="Forbidden: Chat does not belong to user.")
    except Exception as e:
        logger.exception(f"Authentication/Authorization error for user {username} on chat {chat_id}.")
        raise HTTPException(status_code=401, detail="Unauthorized or invalid chat.")


    try:
        # 1. Insert User Message
        db.execute("INSERT INTO messages (chat_id, role, content) VALUES (?, ?, ?)", (chat_id, "user", prompt))
        db.commit()
        logger.info(f"User message added to chat {chat_id} by {username}.")

        # 2. Fetch Chat History (including the new user message)
        messages_raw = db.execute("SELECT role, content FROM messages WHERE chat_id = ? ORDER BY timestamp ASC", (chat_id,)).fetchall()
        messages = [dict(m) for m in messages_raw] # Convert Row objects to dictionaries
        logger.debug(f"Fetched {len(messages)} messages for chat {chat_id}.")

        # 3. Get Chat Summary
        summary_row = db.execute("SELECT summary FROM chat_summaries WHERE chat_id = ?", (chat_id,)).fetchone()
        summary = summary_row["summary"] if summary_row and summary_row["summary"] else "No prior summary available."
        logger.debug(f"Current summary for chat {chat_id}: '{summary[:50]}...'")


        # 4. Process File Content (if any)
        file_text = ""
        if file and file.filename:
            logger.info(f"Processing uploaded file: {file.filename}")
            if file.filename.lower().endswith(".pptx"):
                try:
                    file_bytes = await file.read()
                    prs = Presentation(io.BytesIO(file_bytes))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                file_text += shape.text + "\n"
                    if not file_text.strip():
                        logger.warning(f"PPTX file {file.filename} contained no readable text.")
                        file_text = f"The provided '{file.filename}' contained no extractable text."
                    else:
                         logger.info(f"Extracted {len(file_text)} characters from {file.filename}.")
                except Exception as e:
                    logger.error(f"Error processing PPTX file {file.filename}: {e}")
                    file_text = f"Error: Could not process '{file.filename}'. Details: {e}"
            else:
                logger.warning(f"Unsupported file type uploaded: {file.filename}. Only .pptx is supported.")
                file_text = f"Unsupported file type: '{file.filename}'. Only PowerPoint files (.pptx) are processed."
        elif file and not file.filename: # Handle cases where file might be empty but present
             logger.warning("Received an empty file upload (no filename).")


        # 5. Integrate with FAISS (Vector Store)
        try:
            # Ensure global_embedding_model is not None before passing
            if global_embedding_model is None:
                raise RuntimeError("Embedding model is not initialized. Cannot use FAISS.")

            if chat_id not in faiss_cache:
                faiss_cache[chat_id] = get_faiss_index(chat_id)
            faiss_index = faiss_cache[chat_id]
            if prompt.strip():
                faiss_index.add_documents([Document(page_content=prompt)])
                save_faiss_index(faiss_index, chat_id)
                logger.debug(f"FAISS index updated with current prompt for chat {chat_id}.")

            retrieved_docs = faiss_index.max_marginal_relevance_search(prompt, k=5, fetch_k=15, lambda_mult=0.7)
            retrieved = "\n".join([doc.page_content for doc in retrieved_docs]) if retrieved_docs else "No relevant information retrieved from vector store."
            logger.debug(f"Retrieved {len(retrieved_docs)} docs from FAISS.")
        except Exception as e:
            logger.error(f"Error interacting with FAISS for chat {chat_id}: {e}")
            retrieved = f"Error retrieving context from vector store: {e}"


        # 6. Prepare Recent Chat History
        recent_messages_for_context = messages[-10:] if len(messages) > 10 else messages
        # Exclude the very last user message if we just added it, unless it's the only message.
        # This is to avoid redundancy in the 'Recent Chat' section since 'User: {prompt}' is explicit.
        if len(recent_messages_for_context) > 0 and recent_messages_for_context[-1].get("role") == "user" and recent_messages_for_context[-1].get("content") == prompt:
             recent_messages_for_context = recent_messages_for_context[:-1]

        recent_str = "\n".join([f"{m['role'].capitalize()}: {m['content']}" for m in recent_messages_for_context])
        if not recent_str.strip() and len(messages) > 1: # Only if messages exist beyond the current one
            recent_str = "No recent chat context beyond the current turn."
        elif len(messages) == 1:
            recent_str = "This is the start of the conversation."

        # 7. Construct the Full Context Prompt for LLM
        file_section = f"Content from uploaded file (filename: {file.filename or 'N/A'}):\n{file_text.strip()}" if file_text.strip() else ""

        context_template = PromptTemplate.from_template("""
        Existing Chat Summary:
        {summary}

        Relevant Information:
        {retrieved}

        Recent Chat History:
        {recent}

        {file_section}

        User: {prompt}
        Assistant:""")

         # Temporarily use a high limit to build full context
        temp_limit = MODEL_TOKEN_LIMITS.get(model.lower(), 4096)
        temp_trimmed = trim_context(summary, retrieved, recent_str, file_section, prompt, max_tokens=temp_limit)

        full_context = context_template.format_prompt(
            summary=temp_trimmed["summary"],
            retrieved=temp_trimmed["retrieved"],
            recent=temp_trimmed["recent"],
            file_section=temp_trimmed["file_section"],
            prompt=temp_trimmed["prompt"]
        ).to_string()

        # 🧠 Dynamically calculate response token space AFTER seeing actual context usage
        response_token_budget = calculate_response_token_budget(full_context, model)

        duration = time.time() - start
        logger.info(f"Response generated in {duration:.2f} seconds.")
        response_time_logger.info(f"User: {username} | Chat ID: {chat_id} | Model: {model} | Duration: {duration:.2f} seconds")




        logger.info("------ PROMPT CONTEXT BEGIN ------")
        logger.info(full_context)
        logger.info("------ PROMPT CONTEXT END --------")

        # 8. Get LLM Response
        llm = get_llm(model)
        prompt_template = PromptTemplate.from_template("{context}") # Use {context} as the full prompt
        chain = prompt_template | llm

        try:
            assistant_response = chain.invoke({"context": full_context})
            assistant_response = remove_think_tags(assistant_response)
            if not assistant_response.strip():
                assistant_response = "I'm sorry, I couldn't generate a response. Please try again."
                logger.warning(f"LLM returned an empty response for chat {chat_id}.")
        except Exception as e:
            logger.error(f"Error during LLM inference for chat {chat_id}: {e}")
            assistant_response = f"I apologize, but there was an error generating a response: {e}"
            # Re-raise the HTTP exception from OllamaLLM if it was already caught there
            if isinstance(e, HTTPException):
                raise

        # Check for client disconnection before committing response
        await asyncio.sleep(0.01) # Small pause for async check
        if await request.is_disconnected():
            logger.warning(f"Client disconnected for chat {chat_id} after AI generation. Response not saved.")
            return JSONResponse(status_code=200, content={"response": ""}) # Return empty response if disconnected


        # 9. Insert Assistant Message
        if assistant_response.strip():
            db.execute("INSERT INTO messages (chat_id, role, content) VALUES (?, ?, ?)", (chat_id, "assistant", assistant_response))
            db.commit()
            logger.info(f"Assistant response saved for chat {chat_id}.")
        

        # Trigger background summarization (every 10 messages)
        msg_count = db.execute("SELECT COUNT(*) FROM messages WHERE chat_id = ?", (chat_id,)).fetchone()[0]
        if msg_count >= 10 and msg_count % 10 == 0:
            background_tasks.add_task(summarize_chat, DB_NAME, chat_id, model, messages)

        # Trigger background title generation (after first user + assistant message)
        if msg_count == 2 and messages[0].get("role") == "user":
            background_tasks.add_task(generate_title, DB_NAME, chat_id, messages[0]["content"], model)

        # 10. Periodically Summarize Chat
        # msg_count = db.execute("SELECT COUNT(*) FROM messages WHERE chat_id = ?", (chat_id,)).fetchone()[0]
        # if msg_count >= 10 and msg_count % 10 == 0: # Summarize every 10 messages (adjust as needed)
        #     logger.info(f"Initiating summarization for chat {chat_id} (message count: {msg_count}).")
        #     full_chat_for_summary = "\n".join([
        #         f"{m['role'].capitalize()}: {remove_think_tags(m['content'])}" for m in messages
        #     ]) # Use full history for summary
        #     if full_chat_for_summary.strip():
        #         text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=0) # Adjust chunk_size/overlap if needed
        #         docs = [doc for doc in text_splitter.create_documents([full_chat_for_summary]) if doc.page_content.strip()]
        #         if docs:
        #             try:
        #                 summarize_llm = get_llm(model) # Use the chosen model for summarization
        #                 summary_chain = load_summarize_chain(summarize_llm, chain_type="refine")
    
        #                 new_summary = summary_chain.run(docs).strip()
        #                 if new_summary:
        #                     db.execute("INSERT OR REPLACE INTO chat_summaries (chat_id, summary) VALUES (?, ?)", (chat_id, new_summary))
        #                     db.commit()
        #                     logger.info(f"✅ Summary updated for chat {chat_id}")
        #                 else:
        #                     logger.warning(f"Summarization returned an empty string for chat {chat_id}.")
        #             except Exception as e:
        #                 logger.error(f"❌ Summarization error for chat {chat_id}: {e}")
        #         else:
        #             logger.warning(f"⚠️ Skipping summarization: no valid documents generated from chat history for chat {chat_id}.")
        #     else:
        #         logger.warning(f"⚠️ Skipping summarization: full chat history is empty for chat {chat_id}.")

        # 11. Title Generation for New Chats (after the first user message + assistant response)
        # Check if this is the first assistant response in a new chat
        # The 'messages' list already includes the user's latest prompt.
        # If the total message count is 2 (user + assistant), then it's the first turn.
        # if msg_count == 2 and messages[0].get("role") == "user":
        #     logger.info(f"Attempting title generation for new chat {chat_id}.")
        #     try:
        #         title_template = PromptTemplate.from_template("""
        #         Generate a short, 2–5 word chat title for this user message.

        #         User: {first_user_message}
        #         Only respond with the title. No quotes or punctuation.
        #         """)

        #         title_chain = title_template | OllamaLLM(model=model)

        #         first_user_message_content = messages[0].get("content", "")
        #         if first_user_message_content:
        #             raw_title = remove_think_tags(title_chain.invoke({"first_user_message": first_user_message_content})).strip()
        #             if raw_title:
        #                 # Clean the title: take first line, remove quotes, limit length
        #                 title = raw_title.split("\n")[0].replace('"', '').replace("Title:", "").strip()[:60]
        #                 if not title: # If cleaning resulted in empty string
        #                     title = "New Chat" # Fallback title
        #                 db.execute("UPDATE chats SET title = ? WHERE id = ?", (title, chat_id))
        #                 db.commit()
        #                 logger.info(f"Chat {chat_id} titled: '{title}'.")
        #             else:
        #                 logger.warning(f"Title generation returned empty string for chat {chat_id}.")
        #         else:
        #             logger.warning(f"First user message content was empty for title generation in chat {chat_id}.")
        #     except Exception as e:
        #         logger.error(f"⚠️ Title generation failed for chat {chat_id}: {e}")
        #         # You might want to update the title to a generic one if generation fails
        #         db.execute("UPDATE chats SET title = ? WHERE id = ?", ("Unnamed Chat", chat_id))
        #         db.commit()

        logger.info(f"Response generated in {time.time() - start:.2f} seconds.")
        return JSONResponse(content={"response": assistant_response})
    

    except HTTPException: # Re-raise HTTPExceptions (like from OllamaLLM or auth checks)
        raise
    except Exception as e:
        logger.exception(f"Unhandled error in /api/respond endpoint for user {username}, chat {chat_id}.")
        db.rollback() # Rollback any pending database changes
        return JSONResponse(status_code=500, content={"response": f"❗ Internal error occurred: {e}"})
    


@app.post("/api/delete_chat")
def delete_chat(chat_id: int, request: Request):
    username = request.session.get("user")
    if not username:
        raise HTTPException(status_code=401, detail="Unauthorized")
    db = get_db()
    try:
        # Verify ownership before deleting
        user_id = db.execute("SELECT id FROM users WHERE username = ?", (username,)).fetchone()["id"]
        chat_owner = db.execute("SELECT user_id FROM chats WHERE id = ?", (chat_id,)).fetchone()
        if not chat_owner or chat_owner["user_id"] != user_id:
            logger.warning(f"Unauthorized delete attempt to chat_id {chat_id} by user {username}.")
            raise HTTPException(status_code=403, detail="Forbidden: Chat does not belong to user.")

        db.execute("DELETE FROM messages WHERE chat_id = ?", (chat_id,))
        db.execute("DELETE FROM chat_memory WHERE chat_id = ?", (chat_id,)) # If you use chat_memory table
        db.execute("DELETE FROM chat_summaries WHERE chat_id = ?", (chat_id,))
        db.execute("DELETE FROM chats WHERE id = ?", (chat_id,))
        db.commit()
        logger.info(f"Chat {chat_id} and its associated data deleted by user {username}.")

        # Clean up FAISS index files
        faiss_dir = f"faiss_indexes/chat_{chat_id}"
        if os.path.exists(faiss_dir):
            shutil.rmtree(faiss_dir)
            logger.info(f"Deleted FAISS index directory: {faiss_dir}")

        return {"success": True}
    except Exception as e:
        logger.exception(f"Error deleting chat {chat_id} for user {username}.")
        db.rollback()
        raise HTTPException(status_code=500, detail="Failed to delete chat.")

@app.options("/{rest_of_path:path}")
async def preflight(rest_of_path: str):
    return Response(status_code=204)

@app.on_event("startup")
def startup():
    init_db()


def summarize_chat(db_path, chat_id, model, messages):
    logger.info(f"🎯 Running background summarization for chat {chat_id}")
    try:
        db = sqlite3.connect(db_path)
        db.row_factory = sqlite3.Row
        full_chat = "\n".join([
            f"{m['role'].capitalize()}: {remove_think_tags(m['content'])}"
            for m in messages
        ])
        if full_chat.strip():
            text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=0)
            docs = text_splitter.create_documents([full_chat])
            if docs:
                llm = OllamaLLM(model=model)
                summary_chain = load_summarize_chain(llm, chain_type="refine")
                summary = summary_chain.run(docs).strip()
                if summary:
                    db.execute(
                        "INSERT OR REPLACE INTO chat_summaries (chat_id, summary) VALUES (?, ?)",
                        (chat_id, summary)
                    )
                    db.commit()
    except Exception as e:
        logger.error(f"⚠️ Background summary generation failed: {e}")
    finally:
        db.close()


def generate_title(db_path, chat_id, first_user_message, model):
    logger.info(f"📌 Generating title for chat {chat_id}")
    try:
        db = sqlite3.connect(db_path)
        db.row_factory = sqlite3.Row
        title_template = PromptTemplate.from_template("""
        Generate a short, 2–5 word chat title for this user message.

        User: {first_user_message}
        Only respond with the title. No quotes or punctuation.
        """)
        chain = title_template | OllamaLLM(model=model)
        raw_title = remove_think_tags(chain.invoke({"first_user_message": first_user_message})).strip()
        title = raw_title.split("\n")[0].replace('"', '').replace("Title:", "").strip()[:60]
        if title:
            db.execute("UPDATE chats SET title = ? WHERE id = ?", (title, chat_id))
            db.commit()
    except Exception as e:
        logger.error(f"⚠️ Background title generation failed: {e}")
    finally:
        db.close()
