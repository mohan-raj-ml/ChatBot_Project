import os
import re
import time
import logging
import tiktoken
import requests
import io
from pptx import Presentation
from fastapi import HTTPException
from langchain.prompts import PromptTemplate
from langchain.llms.base import LLM
from langchain.text_splitter import CharacterTextSplitter
from langchain.docstore.document import Document
from langchain.chains.summarize import load_summarize_chain
from databases import Database

# ---- LOGGER ----
logger = logging.getLogger(__name__)
response_time_logger = logging.getLogger("response_time")
response_time_logger.setLevel(logging.INFO)
if not os.path.exists("logs"):
    os.makedirs("logs")
fh = logging.FileHandler("logs/response_times.log")
fh.setFormatter(logging.Formatter('%(asctime)s - %(message)s'))
response_time_logger.addHandler(fh)

# ---- Async PostgreSQL DB ----
DATABASE_URL = "postgresql+asyncpg://postgres:1234@localhost/chatbot_db"

database = Database(DATABASE_URL)

async def get_db():
    return database

# ---- LLM ----
from langchain_core.callbacks import CallbackManagerForLLMRun

from langchain_core.callbacks import CallbackManagerForLLMRun
from langchain.llms.base import LLM
from typing import Optional, List
from fastapi import HTTPException
import requests
import aiohttp
import logging

logger = logging.getLogger(__name__)

import itertools
import threading

OLLAMA_PORTS = [11434, 11436, 11437, 11438]
_round_robin_lock = threading.Lock()
_round_robin_cycle = itertools.cycle(OLLAMA_PORTS)

class OllamaLLM(LLM):
    model: str = "mistral"

    def pick_host(self) -> str:
        with _round_robin_lock:
            port = next(_round_robin_cycle)
        print(f"[DEBUG] Using Ollama on port {port}")
        return f"http://localhost:{port}/api/generate"

    def _call(
        self,
        prompt: str,
        stop: Optional[List[str]] = None,
        run_manager: Optional[CallbackManagerForLLMRun] = None
    ) -> str:
        try:
            url = self.pick_host()
            res = requests.post(url, json={
                "model": self.model,
                "prompt": prompt,
                "stream": False
            }, timeout=300)
            res.raise_for_status()
            return res.json().get("response", "").strip()
        except requests.exceptions.RequestException as e:
            logger.error(f"Error communicating with Ollama server ({self.model}) on {url}: {e}")
            raise HTTPException(status_code=503, detail=f"Failed to connect to Ollama model '{self.model}'.")
        except Exception as e:
            logger.error(f"Unexpected error in OllamaLLM call: {e}")
            raise HTTPException(status_code=500, detail=f"LLM generation failed: {e}")

    async def _acall(
        self,
        prompt: str,
        stop: Optional[List[str]] = None,
        run_manager: Optional[CallbackManagerForLLMRun] = None
    ) -> str:
        try:
            url = self.pick_host()
            async with aiohttp.ClientSession() as session:
                async with session.post(url, json={
                    "model": self.model,
                    "prompt": prompt,
                    "stream": False
                }) as resp:
                    if resp.status == 200:
                        json_resp = await resp.json()
                        return json_resp.get("response", "").strip()
                    else:
                        text = await resp.text()
                        raise HTTPException(status_code=500, detail=f"Ollama async error: {text}")
        except Exception as e:
            logger.error(f"Async error in OllamaLLM ({self.model}): {e}")
            raise HTTPException(status_code=500, detail=f"Ollama async generation failed: {e}")

    @property
    def _identifying_params(self):
        return {"model": self.model}

    @property
    def _llm_type(self):
        return "ollama_custom"


llm_cache = {}
def get_llm(model: str):
    if model not in llm_cache:
        llm_cache[model] = OllamaLLM(model=model)
    return llm_cache[model]

# ---- Tokens ----
MODEL_TOKEN_LIMITS = {
    "gemma:2b": 2048,
    "gemma:7b": 8192,
    "mistral": 8192,
    "mistral:7b": 8192,
    "llama3": 8192,
    "llama3:8b": 8192,
    "llama2": 4096,
    "phi": 2048,
    "tinyllama": 2048,
    "default": 4096
}

def count_tokens(text: str, model: str = "gpt-3.5-turbo") -> int:
    try:
        encoding = tiktoken.encoding_for_model(model)
    except Exception:
        encoding = tiktoken.get_encoding("cl100k_base")
    return len(encoding.encode(text))

def calculate_response_token_budget(full_context: str, model: str) -> int:
    model = model.lower()
    total_limit = next((v for k, v in MODEL_TOKEN_LIMITS.items() if k in model), MODEL_TOKEN_LIMITS["default"])
    used_tokens = count_tokens(full_context, model)
    return max(total_limit - used_tokens - 100, 128)

def trim_context(summary, retrieved, recent, file_section, prompt, max_tokens=3800):
    parts = {
        "prompt": prompt,
        "retrieved": retrieved,
        "summary": summary,
        "recent": recent,
        "file_section": file_section
    }
    used_tokens = 0
    trimmed = {}
    for key, value in parts.items():
        tokens = count_tokens(value)
        if used_tokens + tokens <= max_tokens:
            trimmed[key] = value
            used_tokens += tokens
        else:
            trimmed[key] = f"...[Omitted {key} due to token limit]..."
    return trimmed

# ---- Cleanup ----
def remove_think_tags(text: str) -> str:
    return re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL).strip()

# ---- Summarization ----
async def summarize_chat(db: Database, chat_id: int, model: str, messages: list):
    try:
        full_chat = "\n".join([
            f"{m['role'].capitalize()}: {remove_think_tags(m['content'])}" for m in messages
        ])
        if full_chat.strip():
            text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=0)
            docs = text_splitter.create_documents([full_chat])
            if docs:
                llm = OllamaLLM(model=model)
                summary_chain = load_summarize_chain(llm, chain_type="refine")
                summary = summary_chain.run(docs).strip()
                if summary:
                    await db.execute(
                        "INSERT INTO chat_summaries (chat_id, summary) VALUES (:chat_id, :summary) ON CONFLICT (chat_id) DO UPDATE SET summary = EXCLUDED.summary",
                        {"chat_id": chat_id, "summary": summary}
                    )
    except Exception as e:
        logger.error(f"⚠️ Background summary generation failed: {e}")

# ---- Title Generation ----
async def generate_title(db: Database, chat_id: int, first_user_message: str, model: str):
    try:
        title_template = PromptTemplate.from_template("""
        Generate a short, 2–5 word chat title for this user message.

        User: {first_user_message}
        Only respond with the title. No quotes or punctuation.
        """)
        chain = title_template | OllamaLLM(model=model)
        raw_title = remove_think_tags(chain.invoke({"first_user_message": first_user_message})).strip()
        title = raw_title.split("\n")[0].replace('"', '').replace("Title:", "").strip()[:60]
        if title:
            await db.execute(
                "UPDATE chats SET title = :title WHERE id = :chat_id",
                {"title": title, "chat_id": chat_id}
            )
    except Exception as e:
        logger.error(f"⚠️ Background title generation failed: {e}")

# ---- PPTX Extraction ----
def extract_text_from_pptx(file_bytes: bytes, filename: str) -> str:
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        content = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        return content if content.strip() else f"The provided '{filename}' contained no extractable text."
    except Exception as e:
        logger.error(f"Error processing PPTX file {filename}: {e}")
        return f"Error: Could not process '{filename}'. Details: {e}"

# ---- FAISS Integration ----
from vector_store import get_faiss_index, save_faiss_index, embedding_model as global_embedding_model
faiss_cache = {}

def get_relevant_context(prompt: str, chat_id: int) -> str:
    try:
        if global_embedding_model is None:
            raise RuntimeError("Embedding model is not initialized. Cannot use FAISS.")
        if chat_id not in faiss_cache:
            faiss_cache[chat_id] = get_faiss_index(chat_id)
        faiss_index = faiss_cache[chat_id]
        faiss_index.add_documents([Document(page_content=prompt)])
        save_faiss_index(faiss_index, chat_id)
        retrieved_docs = faiss_index.max_marginal_relevance_search(prompt, k=5, fetch_k=15, lambda_mult=0.7)
        return "\n".join([doc.page_content for doc in retrieved_docs]) if retrieved_docs else "No relevant information retrieved."
    except Exception as e:
        logger.error(f"Error retrieving from FAISS: {e}")
        return f"Error retrieving context from vector store: {e}"
